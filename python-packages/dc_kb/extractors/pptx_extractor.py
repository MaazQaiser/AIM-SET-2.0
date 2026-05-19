from __future__ import annotations

from pathlib import Path

from dc_kb.models import ExtractedDocument, TextChunk


def extract_pptx(path: Path) -> ExtractedDocument:
    ext = path.suffix.lower()
    if ext == ".ppt":
        return _extract_ppt_legacy(path)

    from pptx import Presentation

    prs = Presentation(str(path))
    chunks: list[TextChunk] = []
    for i, slide in enumerate(prs.slides):
        lines = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                lines.append(shape.text.strip())
        text = "\n".join(lines).strip()
        if text:
            chunks.append(TextChunk(text=text, metadata={"slide": i + 1, "source": path.name}))
    return ExtractedDocument(chunks=chunks, metadata={"format": "pptx", "slides": len(prs.slides)})


def _extract_ppt_legacy(path: Path) -> ExtractedDocument:
    try:
        from unstructured.partition.ppt import partition_ppt

        elements = partition_ppt(filename=str(path))
        text = "\n\n".join(str(el) for el in elements).strip()
        if text:
            return ExtractedDocument(
                chunks=[TextChunk(text=text, metadata={"source": path.name})],
                metadata={"format": "ppt"},
            )
    except Exception:
        pass
    raise ValueError(
        f"Legacy .ppt format could not be parsed for {path.name}. Please convert to .pptx and re-upload."
    )
