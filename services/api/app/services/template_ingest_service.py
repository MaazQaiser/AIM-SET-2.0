from __future__ import annotations

import re
import tempfile
from collections import Counter
from html import escape
from io import BytesIO
from pathlib import Path
from typing import Any, Dict, List, Optional

from dc_core.tenancy import TenantContext
from dc_llm.client import LlmClient

from app.config import get_settings
from app.domain.agent_runtime import get_content_generation_runtime
from app.domain.content_studio_repository import get_content_studio_repository
from app.services.office_preview import (
    convert_office_bytes_to_pdf,
    rasterize_presentation_slides,
)

PROMPTS_ROOT = Path(__file__).resolve().parents[4] / "prompts"


def _template_preview_pdf_path(template_id: str, tenant_uuid: str) -> str:
    return f"{tenant_uuid}/{template_id}/preview.pdf"


def load_vision_prompt() -> str:
    path = PROMPTS_ROOT / "content/template_vision/v1.0.0.md"
    if path.is_file():
        return path.read_text(encoding="utf-8")
    return "Convert this slide to HTML/CSS section."


def process_template_ingest(
    ctx: TenantContext,
    template_id: str,
    storage_path: str,
) -> Dict[str, Any]:
    """Rasterize the uploaded source file, persist slide previews, then optionally convert to HTML."""
    repo = get_content_studio_repository()
    settings = get_settings()
    tenant_uuid = repo.template_storage_tenant_uuid(ctx, template_id)
    repo.update_template_progress(
        ctx,
        template_id,
        progress=12,
        stage="reading_source",
        message="Reading uploaded PowerPoint",
    )
    file_bytes = repo.download_template_source(ctx, storage_path)
    ext = Path(storage_path).suffix.lower()

    repo.update_template_progress(
        ctx,
        template_id,
        progress=24,
        stage="rendering_original",
        message="Rendering original slide previews",
    )
    page_images, preview_pdf, preview_pdf_bytes = _prepare_visual_assets(
        repo, tenant_uuid, template_id, file_bytes, ext
    )
    page_count = len(page_images) if page_images else _pdf_page_count(preview_pdf_bytes) or 1
    repo.update_template_progress(
        ctx,
        template_id,
        progress=42,
        stage="extracting_structure",
        message="Extracting slide titles, text, colors, and layouts",
    )
    metadata = _extract_template_metadata(
        file_bytes=file_bytes,
        ext=ext,
        storage_path=storage_path,
        page_count=page_count,
        preview_image_count=len(page_images),
        has_preview_pdf=preview_pdf,
        preview_pdf_bytes=preview_pdf_bytes,
    )
    if not page_images and not preview_pdf:
        repo.finalize_template(
            ctx,
            template_id,
            html="",
            css_variables={},
            page_count=0,
            error="No pages could be rasterized from the uploaded file",
            metadata=_with_processing_metadata(
                metadata,
                progress=100,
                stage="failed",
                message="No pages could be rasterized from the uploaded file",
            ),
        )
        raise ValueError("No pages could be rasterized")

    if page_images:
        repo.save_template_preview_slides(ctx, template_id, page_images)
    repo.update_template_progress(
        ctx,
        template_id,
        progress=58,
        stage="preview_ready",
        message="Original preview is ready",
    )
    thumb = page_images[0] if page_images else None

    try:
        if not page_images:
            raise ValueError("Slide PNGs unavailable; using PDF preview only")

        system = load_vision_prompt()
        runtime = get_content_generation_runtime(ctx)
        client = LlmClient(api_key=settings.llm_api_key or None)
        sections: List[str] = []
        all_vars: Dict[str, str] = {}

        for i, png in enumerate(page_images, start=1):
            repo.update_template_progress(
                ctx,
                template_id,
                progress=64 + round((i - 1) * 24 / max(1, len(page_images))),
                stage="generating_html",
                message=f"Generating HTML/CSS for slide {i} of {len(page_images)}",
            )
            completion = client.complete_vision(
                system=system,
                image_png_bytes=png,
                user_text=f'Slide number: {i}. Emit <section class="slide" data-slide="{i}">...</section>',
                max_tokens=4096,
                model=runtime["model_name"],
                fallback_model=runtime["fallback_model_name"],
            )
            section_html = _clean_section(completion.text, i)
            sections.append(section_html)
            all_vars.update(_extract_css_vars(section_html))

        merged = _merge_template_html(sections, all_vars)
        if _is_placeholder_template_html(merged) and _metadata_has_slide_text(metadata):
            all_vars = _metadata_css_vars(metadata)
            merged = _metadata_to_template_html(metadata, css_variables=all_vars)
            metadata.setdefault("conversion", {})["fallbackGenerated"] = True
        slots = _extract_slots(merged)
        metadata = _with_conversion_metadata(
            metadata,
            html_generated=True,
            section_count=len(sections),
            css_variables=all_vars,
            slots=slots,
        )
        metadata = _with_processing_metadata(
            metadata,
            progress=100,
            stage="ready",
            message="Template HTML/CSS is ready",
        )
        return repo.finalize_template(
            ctx,
            template_id,
            html=merged,
            css_variables=all_vars,
            page_count=page_count,
            thumbnail_bytes=thumb,
            metadata=metadata,
        )
    except Exception as exc:
        fallback_vars = _metadata_css_vars(metadata)
        fallback_html = _metadata_to_template_html(metadata, css_variables=fallback_vars)
        metadata = _with_conversion_metadata(
            metadata,
            html_generated=bool(fallback_html),
            section_count=len(metadata.get("slides") or []),
            css_variables=fallback_vars,
            error=str(exc)[:500],
        )
        metadata.setdefault("conversion", {})["fallbackGenerated"] = bool(fallback_html)
        metadata = _with_processing_metadata(
            metadata,
            progress=100,
            stage="ready",
            message="Original preview and extracted HTML/CSS fallback are ready",
        )
        return repo.finalize_template(
            ctx,
            template_id,
            html=fallback_html,
            css_variables=fallback_vars,
            page_count=page_count,
            thumbnail_bytes=thumb,
            error=str(exc)[:500],
            metadata=metadata,
        )


def ensure_template_preview_slides(ctx: TenantContext, template_id: str) -> int:
    """Rebuild slide PNG previews from the stored source file when missing."""
    repo = get_content_studio_repository()
    row = repo.get_template_row(ctx, template_id)
    storage_path = repo.resolve_template_source_path(ctx, template_id)
    if not storage_path:
        return 0

    slide_count = int((row or {}).get("page_count") or 0)
    if slide_count > 0:
        try:
            repo.download_template_slide(ctx, template_id, 1)
            return slide_count
        except FileNotFoundError:
            pass

    tenant_uuid = repo.template_storage_tenant_uuid(ctx, template_id)
    file_bytes = repo.download_template_source(ctx, storage_path)
    ext = Path(storage_path).suffix.lower()
    page_images, _, _ = _prepare_visual_assets(repo, tenant_uuid, template_id, file_bytes, ext)
    if not page_images:
        return 1 if repo.get_template_preview_pdf(ctx, template_id) else 0

    saved = repo.save_template_preview_slides(ctx, template_id, page_images)
    return saved


def ensure_template_metadata(ctx: TenantContext, template_id: str) -> Dict[str, Any]:
    """Populate slide structure metadata for old templates that still have an empty analysis."""
    repo = get_content_studio_repository()
    row = repo.get_template_row(ctx, template_id) or {}
    current = row.get("metadata") or {}
    if isinstance(current, dict) and current.get("slides"):
        return current

    storage_path = repo.resolve_template_source_path(ctx, template_id)
    if not storage_path:
        return current if isinstance(current, dict) else {}

    file_bytes = repo.download_template_source(ctx, storage_path)
    ext = Path(storage_path).suffix.lower()
    page_count = int(row.get("page_count") or 0) or 1
    preview_pdf_bytes: Optional[bytes] = None
    has_preview_pdf = False
    if ext == ".pdf":
        preview_pdf_bytes = file_bytes
        has_preview_pdf = True
    elif ext == ".ppt":
        preview_pdf_bytes = repo.get_template_preview_pdf(ctx, template_id)
        has_preview_pdf = bool(preview_pdf_bytes)
        if not preview_pdf_bytes:
            try:
                preview_pdf_bytes = convert_office_bytes_to_pdf(file_bytes, ext, allow_text_fallback=True)
                has_preview_pdf = True
            except Exception:
                preview_pdf_bytes = None

    metadata = _extract_template_metadata(
        file_bytes=file_bytes,
        ext=ext,
        storage_path=storage_path,
        page_count=page_count,
        preview_image_count=page_count,
        has_preview_pdf=has_preview_pdf,
        preview_pdf_bytes=preview_pdf_bytes,
    )
    css_vars = row.get("css_variables") if isinstance(row.get("css_variables"), dict) else {}
    has_html = bool(row.get("html"))
    metadata = _with_conversion_metadata(
        metadata,
        html_generated=has_html,
        section_count=_count_html_slides(str(row.get("html") or "")),
        css_variables=css_vars,
    )
    metadata = _with_processing_metadata(
        metadata,
        progress=100,
        stage="ready",
        message="Slide structure extracted",
    )

    patch: Dict[str, Any] = {"metadata": metadata}
    if (not has_html or _is_placeholder_template_html(str(row.get("html") or ""))) and _metadata_has_slide_text(metadata):
        fallback_vars = _metadata_css_vars(metadata)
        fallback_html = _metadata_to_template_html(metadata, css_variables=fallback_vars)
        if fallback_html:
            metadata = _with_conversion_metadata(
                metadata,
                html_generated=True,
                section_count=len(metadata.get("slides") or []),
                css_variables=fallback_vars,
            )
            metadata.setdefault("conversion", {})["fallbackGenerated"] = True
            patch.update(
                {
                    "html": fallback_html,
                    "cssVariables": fallback_vars,
                    "pageCount": len(metadata.get("slides") or []) or page_count,
                    "metadata": metadata,
                }
            )
    repo.update_template(ctx, template_id, patch)
    return metadata


def _prepare_visual_assets(
    repo,
    tenant_uuid: str,
    template_id: str,
    file_bytes: bytes,
    ext: str,
) -> tuple[List[bytes], bool, Optional[bytes]]:
    preview_pdf = False
    preview_pdf_bytes: Optional[bytes] = None
    page_images: List[bytes] = []

    if ext in {".png", ".jpg", ".jpeg"}:
        return [file_bytes], False, None

    if ext in {".ppt", ".pptx"}:
        try:
            pdf_bytes = convert_office_bytes_to_pdf(file_bytes, ext, allow_text_fallback=False)
            preview_pdf_bytes = pdf_bytes
            repo.upload_template_blob(
                _template_preview_pdf_path(template_id, tenant_uuid),
                pdf_bytes,
                content_type="application/pdf",
            )
            preview_pdf = True
        except Exception:
            preview_pdf = False

        try:
            page_images = rasterize_presentation_slides(file_bytes, ext)
            return page_images, preview_pdf, preview_pdf_bytes
        except Exception:
            if preview_pdf:
                return [], True, preview_pdf_bytes
            raise

    if ext == ".pdf":
        repo.upload_template_blob(
            _template_preview_pdf_path(template_id, tenant_uuid),
            file_bytes,
            content_type="application/pdf",
        )
        return _pdf_to_pngs_from_bytes(file_bytes), True, file_bytes

    raise ValueError(f"Unsupported extension: {ext}")


def _extract_template_metadata(
    *,
    file_bytes: bytes,
    ext: str,
    storage_path: str,
    page_count: int,
    preview_image_count: int,
    has_preview_pdf: bool,
    preview_pdf_bytes: Optional[bytes],
) -> Dict[str, Any]:
    """Best-effort source understanding for preview and future generation planning."""
    source_file_name = Path(storage_path).name
    metadata: Dict[str, Any] = {
        "source": {
            "fileName": source_file_name,
            "extension": ext.lstrip("."),
        },
        "slideCount": page_count,
        "slides": _fallback_slides(page_count),
        "design": {
            "colors": [],
            "fonts": [],
            "layouts": [],
        },
        "conversion": {
            "previewImageCount": preview_image_count,
            "hasPreviewPdf": has_preview_pdf,
            "htmlGenerated": False,
            "sectionCount": 0,
            "cssVariables": [],
        },
    }

    extracted: Dict[str, Any] = {}
    try:
        if ext == ".pptx":
            extracted = _extract_pptx_metadata(file_bytes)
        elif ext == ".pdf" and preview_pdf_bytes:
            extracted = _extract_pdf_metadata(preview_pdf_bytes)
        elif ext == ".ppt" and preview_pdf_bytes:
            extracted = _extract_pdf_metadata(preview_pdf_bytes)
            extracted.setdefault("sourceFormatNote", "Legacy PPT text was extracted from the converted PDF preview.")
    except Exception as exc:
        metadata["extractionError"] = str(exc)[:300]

    if extracted:
        metadata = _merge_metadata(metadata, extracted)

    metadata["slideCount"] = max(
        int(metadata.get("slideCount") or 0),
        page_count,
        len(metadata.get("slides") or []),
    )
    metadata["slides"] = _normalize_slides(metadata.get("slides") or [], int(metadata["slideCount"]))
    metadata["design"] = _summarize_design(metadata["slides"], metadata.get("design") or {})
    return metadata


def _extract_pptx_metadata(file_bytes: bytes) -> Dict[str, Any]:
    from pptx import Presentation  # type: ignore

    prs = Presentation(BytesIO(file_bytes))
    slides: List[Dict[str, Any]] = []
    all_colors: Counter[str] = Counter()
    all_fonts: Counter[str] = Counter()
    layouts: Counter[str] = Counter()

    for index, slide in enumerate(prs.slides, start=1):
        layout = str(getattr(slide.slide_layout, "name", "") or f"Layout {index}")
        layouts[layout] += 1
        text_blocks: List[str] = []
        slide_colors: Counter[str] = Counter()
        slide_fonts: Counter[str] = Counter()
        image_count = 0
        table_count = 0
        chart_count = 0

        background = _fill_color(getattr(getattr(slide, "background", None), "fill", None))
        if background:
            slide_colors[background] += 1
            all_colors[background] += 1

        for shape in slide.shapes:
            if _shape_has_image(shape):
                image_count += 1
            if bool(getattr(shape, "has_table", False)):
                table_count += 1
            if bool(getattr(shape, "has_chart", False)):
                chart_count += 1

            for color in _shape_colors(shape):
                slide_colors[color] += 1
                all_colors[color] += 1

            if bool(getattr(shape, "has_text_frame", False)):
                text = _shape_text(shape)
                if text:
                    text_blocks.append(text)
                for font in _shape_fonts(shape):
                    slide_fonts[font] += 1
                    all_fonts[font] += 1

        title = _slide_title(slide, text_blocks, index)
        slides.append(
            {
                "slide": index,
                "name": title,
                "title": title,
                "layout": layout,
                "text": _clamp_text("\n".join(text_blocks), 6000),
                "textBlocks": [_clamp_text(block, 1200) for block in text_blocks[:24]],
                "colors": _counter_keys(slide_colors),
                "fonts": _counter_keys(slide_fonts),
                "shapeCount": len(slide.shapes),
                "imageCount": image_count,
                "tableCount": table_count,
                "chartCount": chart_count,
            }
        )

    width = int(getattr(prs, "slide_width", 0) or 0)
    height = int(getattr(prs, "slide_height", 0) or 0)
    return {
        "slideCount": len(slides),
        "slides": slides,
        "design": {
            "colors": _counter_keys(all_colors),
            "fonts": _counter_keys(all_fonts),
            "layouts": _counter_keys(layouts, limit=12),
            "slideSize": {"widthEmu": width, "heightEmu": height},
        },
    }


def _extract_pdf_metadata(pdf_bytes: bytes) -> Dict[str, Any]:
    import fitz  # type: ignore

    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    slides: List[Dict[str, Any]] = []
    all_colors: Counter[str] = Counter()
    all_fonts: Counter[str] = Counter()
    try:
        for index, page in enumerate(doc, start=1):
            text_blocks: List[str] = []
            page_colors: Counter[str] = Counter()
            page_fonts: Counter[str] = Counter()
            data = page.get_text("dict") or {}
            for block in data.get("blocks") or []:
                lines: List[str] = []
                for line in block.get("lines") or []:
                    line_text = ""
                    for span in line.get("spans") or []:
                        line_text += str(span.get("text") or "")
                        font = str(span.get("font") or "").strip()
                        if font:
                            page_fonts[font] += 1
                            all_fonts[font] += 1
                        color = _int_color_to_hex(span.get("color"))
                        if color:
                            page_colors[color] += 1
                            all_colors[color] += 1
                    if line_text.strip():
                        lines.append(line_text.strip())
                if lines:
                    text_blocks.append(" ".join(lines))

            page_text = "\n".join(text_blocks) or str(page.get_text("text") or "").strip()
            title = _first_text_line(page_text) or f"Slide {index}"
            slides.append(
                {
                    "slide": index,
                    "name": title,
                    "title": title,
                    "layout": "PDF page",
                    "text": _clamp_text(page_text, 6000),
                    "textBlocks": [_clamp_text(block, 1200) for block in text_blocks[:24]],
                    "colors": _counter_keys(page_colors),
                    "fonts": _counter_keys(page_fonts),
                    "shapeCount": 0,
                    "imageCount": len([b for b in data.get("blocks") or [] if b.get("type") == 1]),
                    "tableCount": 0,
                    "chartCount": 0,
                }
            )
    finally:
        doc.close()

    return {
        "slideCount": len(slides),
        "slides": slides,
        "design": {
            "colors": _counter_keys(all_colors),
            "fonts": _counter_keys(all_fonts),
            "layouts": ["PDF page"] if slides else [],
        },
    }


def _with_conversion_metadata(
    metadata: Dict[str, Any],
    *,
    html_generated: bool,
    section_count: int,
    css_variables: Dict[str, str],
    error: Optional[str] = None,
    slots: Optional[List[Dict[str, str]]] = None,
) -> Dict[str, Any]:
    out = dict(metadata)
    conversion = dict(out.get("conversion") or {})
    conversion.update(
        {
            "htmlGenerated": html_generated,
            "sectionCount": section_count,
            "cssVariables": sorted(css_variables.keys()),
        }
    )
    if slots is not None:
        conversion["slots"] = slots
    if error:
        conversion["error"] = error
    out["conversion"] = conversion
    return out


def _with_processing_metadata(
    metadata: Dict[str, Any],
    *,
    progress: int,
    stage: str,
    message: str,
) -> Dict[str, Any]:
    out = dict(metadata)
    out["processing"] = {
        "progress": max(0, min(100, int(progress))),
        "stage": stage,
        "message": message,
    }
    return out


def _metadata_has_slide_text(metadata: Dict[str, Any]) -> bool:
    for slide in metadata.get("slides") or []:
        if str(slide.get("title") or "").strip() and str(slide.get("title")) != f"Slide {slide.get('slide')}":
            return True
        if str(slide.get("text") or "").strip():
            return True
    return False


def _metadata_css_vars(metadata: Dict[str, Any]) -> Dict[str, str]:
    colors = ((metadata.get("design") or {}).get("colors") or []) if isinstance(metadata.get("design"), dict) else []
    primary = str(colors[0]) if colors else "#2563eb"
    accent = str(colors[1]) if len(colors) > 1 else primary
    return {
        "--template-primary": primary,
        "--template-accent": accent,
        "--template-bg": "#ffffff",
        "--template-text": "#111827",
        "--template-muted": "#6b7280",
    }


def _metadata_to_template_html(metadata: Dict[str, Any], *, css_variables: Dict[str, str]) -> str:
    slides = metadata.get("slides") or []
    if not slides:
        return ""
    root_vars = "\n".join(f"  {key}: {value};" for key, value in css_variables.items())
    sections: List[str] = []
    for slide in slides:
        slide_no = int(slide.get("slide") or len(sections) + 1)
        title = str(slide.get("title") or slide.get("name") or f"Slide {slide_no}").strip()
        layout = str(slide.get("layout") or "").strip()
        text_blocks = [str(block).strip() for block in (slide.get("textBlocks") or []) if str(block).strip()]
        if not text_blocks and str(slide.get("text") or "").strip():
            text_blocks = [line.strip() for line in str(slide.get("text")).splitlines() if line.strip()]
        # Skip the title if it already appears as the first text block to avoid duplication.
        body_blocks = [b for b in text_blocks[:8] if b.lower() != title.lower()]
        body = "\n".join(f"<p>{escape(block)}</p>" for block in body_blocks)
        color_chips = "\n".join(
            f'<span class="color-chip" style="background:{escape(str(color))}"></span>'
            for color in (slide.get("colors") or [])[:8]
            if str(color).strip()
        )
        sections.append(
            "\n".join(
                [
                    f'<section class="slide" data-slide="{slide_no}">',
                    f"<h1>{escape(title)}</h1>",
                    f'<p class="slide-layout">{escape(layout)}</p>' if layout else "",
                    f'<div class="slide-body">{body}</div>' if body else "",
                    f'<div class="slide-colors">{color_chips}</div>' if color_chips else "",
                    f'<div class="slide-number">{slide_no}</div>',
                    "</section>",
                ]
            )
        )
    return (
        '<!DOCTYPE html>'
        '<html lang="en">'
        '<head>'
        '<meta charset="utf-8" />'
        '<meta name="viewport" content="width=1280, initial-scale=1" />'
        '<link rel="preconnect" href="https://fonts.googleapis.com" />'
        '<link rel="preconnect" href="https://fonts.gstatic.com" crossorigin />'
        '<link href="https://fonts.googleapis.com/css2?family=Urbanist:wght@300;400;500;600;700;800&display=swap" rel="stylesheet" />'
        "<style>"
        f":root {{\n{root_vars}\n}}\n"
        "*, *::before, *::after { box-sizing: border-box; }"
        "html, body { margin: 0; padding: 0; background: #1c1c1e; }"
        # Slide canvas: 1280×720, scaled to fit via viewport meta.
        # Background is driven by the extracted --template-primary variable so it
        # matches the actual deck colors rather than a hardcoded blue.
        ".slide {"
        "  position: relative;"
        "  width: 1280px;"
        "  height: 720px;"
        "  overflow: hidden;"
        "  margin: 0 auto;"
        "  padding: 80px 96px;"
        "  font-family: 'Urbanist', Arial, sans-serif;"
        "  color: var(--template-text, #111827);"
        "  background: var(--template-bg, var(--template-primary, #ffffff));"
        "}"
        ".slide h1 {"
        "  margin: 0 0 20px;"
        "  font-size: 64px;"
        "  font-weight: 800;"
        "  line-height: 1.05;"
        "  letter-spacing: -0.02em;"
        "  color: #ffffff;"
        "  max-width: 840px;"
        "}"
        ".slide-layout {"
        "  display: inline-block;"
        "  margin: 0 0 24px;"
        "  padding: 4px 12px;"
        "  border: 1px solid rgba(255,255,255,.35);"
        "  border-radius: 4px;"
        "  font-size: 13px;"
        "  font-weight: 500;"
        "  color: rgba(255,255,255,.7);"
        "  letter-spacing: 0.06em;"
        "  text-transform: uppercase;"
        "}"
        ".slide-body {"
        "  max-width: 800px;"
        "  display: grid;"
        "  gap: 14px;"
        "  font-size: 22px;"
        "  line-height: 1.55;"
        "  color: rgba(255,255,255,.85);"
        "}"
        ".slide-body p { margin: 0; }"
        ".slide-colors {"
        "  position: absolute;"
        "  left: 96px;"
        "  bottom: 56px;"
        "  display: flex;"
        "  gap: 8px;"
        "}"
        ".color-chip {"
        "  width: 24px;"
        "  height: 24px;"
        "  border-radius: 50%;"
        "  border: 2px solid rgba(255,255,255,.3);"
        "}"
        ".slide-number {"
        "  position: absolute;"
        "  right: 48px;"
        "  bottom: 40px;"
        "  color: rgba(255,255,255,.4);"
        "  font-size: 16px;"
        "  font-weight: 500;"
        "}"
        "</style>"
        "</head>"
        "<body>"
        + "\n".join(sections)
        + "</body></html>"
    )


def _is_placeholder_template_html(html: str) -> bool:
    visible = re.sub(r"<style[\s\S]*?</style>", "", html, flags=re.IGNORECASE)
    visible = re.sub(r"<[^>]+>", " ", visible)
    words = [word.lower() for word in re.findall(r"[A-Za-z0-9]+", visible)]
    if not words:
        return True
    meaningful = [word for word in words if word not in {"slide", "template"} and not word.isdigit()]
    return len(meaningful) <= 2 and len(words) <= 30


def _count_html_slides(html: str) -> int:
    if not html:
        return 0
    sections = re.findall(r"<section\b", html, flags=re.IGNORECASE)
    if sections:
        return len(sections)
    return len(re.findall(r"data-slide=", html, flags=re.IGNORECASE))


def _merge_metadata(base: Dict[str, Any], extracted: Dict[str, Any]) -> Dict[str, Any]:
    merged = dict(base)
    for key, value in extracted.items():
        if key == "design":
            design = dict(merged.get("design") or {})
            design.update(value or {})
            merged["design"] = design
        elif key == "slides":
            merged["slides"] = value or merged.get("slides") or []
        else:
            merged[key] = value
    return merged


def _fallback_slides(page_count: int) -> List[Dict[str, Any]]:
    return [
        {
            "slide": i,
            "name": f"Slide {i}",
            "title": f"Slide {i}",
            "layout": "",
            "text": "",
            "textBlocks": [],
            "colors": [],
            "fonts": [],
            "shapeCount": 0,
            "imageCount": 0,
            "tableCount": 0,
            "chartCount": 0,
        }
        for i in range(1, max(0, page_count) + 1)
    ]


def _normalize_slides(slides: List[Dict[str, Any]], slide_count: int) -> List[Dict[str, Any]]:
    by_index = {int(s.get("slide") or idx + 1): s for idx, s in enumerate(slides) if isinstance(s, dict)}
    normalized: List[Dict[str, Any]] = []
    for index in range(1, max(0, slide_count) + 1):
        slide = dict(by_index.get(index) or {})
        title = str(slide.get("title") or slide.get("name") or f"Slide {index}").strip()
        normalized.append(
            {
                "slide": index,
                "name": str(slide.get("name") or title),
                "title": title,
                "layout": str(slide.get("layout") or ""),
                "text": str(slide.get("text") or ""),
                "textBlocks": [str(v) for v in (slide.get("textBlocks") or []) if str(v).strip()],
                "colors": _unique_strings(slide.get("colors") or [], limit=10),
                "fonts": _unique_strings(slide.get("fonts") or [], limit=10),
                "shapeCount": int(slide.get("shapeCount") or 0),
                "imageCount": int(slide.get("imageCount") or 0),
                "tableCount": int(slide.get("tableCount") or 0),
                "chartCount": int(slide.get("chartCount") or 0),
            }
        )
    return normalized


def _summarize_design(slides: List[Dict[str, Any]], design: Dict[str, Any]) -> Dict[str, Any]:
    colors = Counter()
    fonts = Counter()
    layouts = Counter()
    for slide in slides:
        for color in slide.get("colors") or []:
            colors[str(color)] += 1
        for font in slide.get("fonts") or []:
            fonts[str(font)] += 1
        layout = str(slide.get("layout") or "").strip()
        if layout:
            layouts[layout] += 1

    out = dict(design)
    out["colors"] = _unique_strings(out.get("colors") or _counter_keys(colors), limit=12)
    out["fonts"] = _unique_strings(out.get("fonts") or _counter_keys(fonts), limit=12)
    out["layouts"] = _unique_strings(out.get("layouts") or _counter_keys(layouts), limit=12)
    return out


def _slide_title(slide: Any, text_blocks: List[str], slide_index: int) -> str:
    try:
        title_shape = slide.shapes.title
        title = _shape_text(title_shape)
        if title:
            return _first_text_line(title) or title
    except Exception:
        pass
    return _first_text_line("\n".join(text_blocks)) or f"Slide {slide_index}"


def _shape_text(shape: Any) -> str:
    try:
        value = str(getattr(shape, "text", "") or "").strip()
        if value:
            return re.sub(r"\n{3,}", "\n\n", value)
    except Exception:
        return ""
    return ""


def _shape_fonts(shape: Any) -> List[str]:
    fonts: List[str] = []
    try:
        text_frame = shape.text_frame
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                name = str(getattr(run.font, "name", "") or "").strip()
                if name:
                    fonts.append(name)
            paragraph_font = str(getattr(paragraph.font, "name", "") or "").strip()
            if paragraph_font:
                fonts.append(paragraph_font)
    except Exception:
        pass
    return _unique_strings(fonts, limit=16)


def _shape_colors(shape: Any) -> List[str]:
    colors: List[str] = []
    fill = _fill_color(getattr(shape, "fill", None))
    line = _line_color(getattr(shape, "line", None))
    if fill:
        colors.append(fill)
    if line:
        colors.append(line)
    try:
        if bool(getattr(shape, "has_text_frame", False)):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    color = _font_color(getattr(run, "font", None))
                    if color:
                        colors.append(color)
    except Exception:
        pass
    return _unique_strings(colors, limit=16)


def _shape_has_image(shape: Any) -> bool:
    try:
        return bool(getattr(shape, "image", None))
    except Exception:
        return False


def _fill_color(fill: Any) -> Optional[str]:
    try:
        return _rgb_to_hex(getattr(getattr(fill, "fore_color", None), "rgb", None))
    except Exception:
        return None


def _line_color(line: Any) -> Optional[str]:
    try:
        return _rgb_to_hex(getattr(getattr(line, "color", None), "rgb", None))
    except Exception:
        return None


def _font_color(font: Any) -> Optional[str]:
    try:
        return _rgb_to_hex(getattr(getattr(font, "color", None), "rgb", None))
    except Exception:
        return None


def _rgb_to_hex(value: Any) -> Optional[str]:
    if value is None:
        return None
    raw = str(value).strip()
    if re.fullmatch(r"[0-9A-Fa-f]{6}", raw):
        return f"#{raw.upper()}"
    if raw.startswith("#") and re.fullmatch(r"#[0-9A-Fa-f]{6}", raw):
        return raw.upper()
    return None


def _int_color_to_hex(value: Any) -> Optional[str]:
    if value is None:
        return None
    try:
        return f"#{int(value) & 0xFFFFFF:06X}"
    except Exception:
        return None


def _counter_keys(counter: Counter[str], limit: int = 10) -> List[str]:
    return [key for key, _count in counter.most_common(limit) if str(key).strip()]


def _unique_strings(values: Any, limit: int = 10) -> List[str]:
    seen = set()
    out: List[str] = []
    for value in values or []:
        text = str(value).strip()
        if not text or text in seen:
            continue
        seen.add(text)
        out.append(text)
        if len(out) >= limit:
            break
    return out


def _first_text_line(text: str) -> str:
    for line in str(text or "").splitlines():
        cleaned = re.sub(r"\s+", " ", line).strip()
        if cleaned:
            return _clamp_text(cleaned, 120)
    return ""


def _clamp_text(text: str, limit: int) -> str:
    normalized = str(text or "").strip()
    if len(normalized) <= limit:
        return normalized
    return normalized[: limit - 3].rstrip() + "..."


def _pdf_page_count(pdf_bytes: Optional[bytes]) -> int:
    if not pdf_bytes:
        return 0
    try:
        import fitz  # type: ignore

        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        try:
            return len(doc)
        finally:
            doc.close()
    except Exception:
        return 0


def _pdf_to_pngs_from_bytes(file_bytes: bytes) -> List[bytes]:
    with tempfile.TemporaryDirectory() as tmp:
        pdf_path = Path(tmp) / "source.pdf"
        pdf_path.write_bytes(file_bytes)
        return _pdf_to_pngs(pdf_path)


def _pdf_to_pngs(pdf_path: Path) -> List[bytes]:
    try:
        from pdf2image import convert_from_path  # type: ignore

        images = convert_from_path(str(pdf_path), dpi=150)
        out: List[bytes] = []
        for img in images:
            import io

            buf = io.BytesIO()
            img.save(buf, format="PNG")
            out.append(buf.getvalue())
        return out
    except Exception:
        pass

    try:
        import fitz  # type: ignore

        doc = fitz.open(str(pdf_path))
        out = []
        for page in doc:
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
            out.append(pix.tobytes("png"))
        doc.close()
        return out
    except Exception as exc:
        raise ValueError(f"PDF rasterize failed: {exc}") from exc


def _clean_section(text: str, slide_num: int) -> str:
    text = text.strip()
    if text.startswith("```"):
        text = re.sub(r"^```[a-z]*\n?", "", text)
        text = re.sub(r"\n?```$", "", text)
    if "<section" not in text.lower():
        text = f'<section class="slide" data-slide="{slide_num}">{text}</section>'
    return text


def _extract_css_vars(html: str) -> Dict[str, str]:
    vars_found: Dict[str, str] = {}
    for m in re.finditer(r"--([a-zA-Z0-9\-]+)\s*:\s*([^;]+);", html):
        vars_found[f"--{m.group(1)}"] = m.group(2).strip()
    return vars_found


def _extract_slots(html: str) -> List[Dict[str, str]]:
    """Return ordered list of {id, type, label} for every data-slot in the HTML."""
    slots: List[Dict[str, str]] = []
    seen: set = set()
    # Match data-slot attributes with optional data-slot-type / data-slot-label on the same tag.
    tag_pattern = re.compile(r"<[^>]+>", re.DOTALL)
    for tag_match in tag_pattern.finditer(html):
        tag = tag_match.group(0)
        slot_id_m = re.search(r'data-slot=["\']([^"\']+)["\']', tag)
        if not slot_id_m:
            continue
        slot_id = slot_id_m.group(1)
        if slot_id in seen:
            continue
        seen.add(slot_id)
        slot_type_m = re.search(r'data-slot-type=["\']([^"\']+)["\']', tag)
        slot_label_m = re.search(r'data-slot-label=["\']([^"\']+)["\']', tag)
        slots.append(
            {
                "id": slot_id,
                "type": slot_type_m.group(1) if slot_type_m else "text",
                "label": slot_label_m.group(1) if slot_label_m else slot_id,
            }
        )
    return slots


def _merge_template_html(sections: List[str], css_variables: Dict[str, str]) -> str:
    root_vars = "\n".join(f"  {k}: {v};" for k, v in css_variables.items())
    body = "\n".join(sections)
    return (
        '<!DOCTYPE html>'
        '<html lang="en">'
        '<head>'
        '<meta charset="utf-8" />'
        # Treat the iframe as 1280px wide so absolute pixel sizes in generated slides
        # render proportionally and scale down automatically on smaller viewports.
        '<meta name="viewport" content="width=1280, initial-scale=1" />'
        '<link rel="preconnect" href="https://fonts.googleapis.com" />'
        '<link rel="preconnect" href="https://fonts.gstatic.com" crossorigin />'
        '<link href="https://fonts.googleapis.com/css2?family=Urbanist:wght@300;400;500;600;700;800&display=swap" rel="stylesheet" />'
        "<style>"
        f":root {{\n{root_vars}\n}}\n"
        # Reset browser defaults so nothing bleeds into the slide canvas.
        "*, *::before, *::after { box-sizing: border-box; }"
        # Neutral dark background so the slide canvas itself (which carries its own background) floats cleanly.
        "html, body { margin: 0; padding: 0; background: #1c1c1e; }"
        # Each section fills the 1280-wide viewport exactly (16:9 = 720px tall).
        # The viewport meta makes the iframe browser scale this down to fit smaller containers.
        "section.slide {"
        "  display: block;"
        "  position: relative;"
        "  width: 1280px;"
        "  height: 720px;"
        "  overflow: hidden;"
        "  margin: 0 auto;"
        "  font-family: var(--body-font, 'Urbanist', Arial, sans-serif);"
        "}"
        # Utility: hide placeholder outlines in preview.
        ".placeholder { display: block; }"
        "</style>"
        "</head>"
        f"<body>{body}</body>"
        "</html>"
    )
