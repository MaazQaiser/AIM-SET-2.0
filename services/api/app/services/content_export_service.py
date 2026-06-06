from __future__ import annotations

import io
import csv
import html as html_lib
import re
import tempfile
import zipfile
from pathlib import Path
from typing import List, Tuple

from dc_core.tenancy import TenantContext

from app.domain.content_studio_repository import get_content_studio_repository


def export_revision(
    ctx: TenantContext,
    revision_id: str,
    fmt: str,
) -> dict:
    if fmt not in ("pdf", "png", "pptx"):
        raise ValueError(f"Unsupported format: {fmt}")

    repo = get_content_studio_repository()
    revision = repo.get_revision(ctx, revision_id)
    if not revision:
        raise ValueError(f"Revision not found: {revision_id}")

    html = revision["html"]
    data = render_html_export(html, fmt)

    result = repo.create_export(ctx, revision_id, fmt=fmt, file_bytes=data)
    project_id = revision.get("projectId")
    if project_id:
        repo.update_project(ctx, str(project_id), {"status": "exported"})
    return result


def render_html_export(html: str, fmt: str) -> bytes:
    if fmt not in ("pdf", "png", "pptx", "csv"):
        raise ValueError(f"Unsupported format: {fmt}")
    if fmt == "pdf":
        return _export_pdf(html)
    if fmt == "png":
        return _export_png_zip(html)
    if fmt == "csv":
        return _export_csv(html)
    return _export_pptx(html)


def render_html_png(html: str) -> bytes:
    sections = _split_sections(html)
    section = sections[0] if sections else html
    try:
        from playwright.sync_api import sync_playwright  # type: ignore

        with sync_playwright() as p:
            browser = p.chromium.launch(headless=True)
            page = browser.new_page(viewport={"width": 1280, "height": 720})
            page.set_content(_wrap_section_page(section), wait_until="networkidle")
            png = page.screenshot(full_page=False)
            browser.close()
            return png
    except Exception:
        return _export_png_fallback([section])[0][1]


def export_revision_file_bytes(ctx: TenantContext, revision_id: str, fmt: str) -> bytes:
    if fmt not in ("pdf", "png", "pptx"):
        raise ValueError(f"Unsupported format: {fmt}")

    repo = get_content_studio_repository()
    revision = repo.get_revision(ctx, revision_id)
    if not revision:
        raise ValueError(f"Revision not found: {revision_id}")

    html = revision["html"]
    if fmt == "pdf":
        return _export_pdf(html)
    if fmt == "png":
        return _export_png_zip(html)
    return _export_pptx(html)


def _export_pdf(html: str) -> bytes:
    try:
        from playwright.sync_api import sync_playwright  # type: ignore

        with sync_playwright() as p:
            browser = p.chromium.launch(headless=True)
            page = browser.new_page(viewport={"width": 1280, "height": 720})
            page.set_content(html, wait_until="networkidle")
            pdf = page.pdf(
                width="1280px",
                height="720px",
                print_background=True,
                prefer_css_page_size=True,
            )
            browser.close()
            return pdf
    except Exception:
        return _export_pdf_fallback(html)


def _export_pdf_fallback(html: str) -> bytes:
    import fitz  # type: ignore

    doc = fitz.open()
    for _i, section_html in enumerate(_split_sections(html), start=1):
        page = doc.new_page(width=1280, height=720)
        page.insert_htmlbox(fitz.Rect(0, 0, 1280, 720), section_html)
    buf = io.BytesIO()
    doc.save(buf)
    doc.close()
    return buf.getvalue()


def _export_png_zip(html: str) -> bytes:
    sections = _split_sections(html)
    if not sections:
        sections = [html]

    images: List[Tuple[str, bytes]] = []
    try:
        from playwright.sync_api import sync_playwright  # type: ignore

        with sync_playwright() as p:
            browser = p.chromium.launch(headless=True)
            for i, section in enumerate(sections, start=1):
                page = browser.new_page(viewport={"width": 1280, "height": 720})
                wrapped = _wrap_section_page(section)
                page.set_content(wrapped, wait_until="networkidle")
                png = page.screenshot(full_page=False)
                images.append((f"slide-{i}.png", png))
            browser.close()
    except Exception:
        images = _export_png_fallback(sections)

    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for name, data in images:
            zf.writestr(name, data)
    return buf.getvalue()


def _export_png_fallback(sections: List[str]) -> List[Tuple[str, bytes]]:
    import fitz  # type: ignore

    out: List[Tuple[str, bytes]] = []
    for i, section in enumerate(sections, start=1):
        doc = fitz.open()
        page = doc.new_page(width=1280, height=720)
        page.insert_htmlbox(fitz.Rect(0, 0, 1280, 720), section)
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        out.append((f"slide-{i}.png", pix.tobytes("png")))
        doc.close()
    return out


def _export_pptx(html: str) -> bytes:
    from pptx import Presentation  # type: ignore
    from pptx.util import Inches

    sections = _split_sections(html)
    if not sections:
        sections = [html]

    pngs: List[bytes] = []
    try:
        from playwright.sync_api import sync_playwright  # type: ignore

        with sync_playwright() as p:
            browser = p.chromium.launch(headless=True)
            for section in sections:
                page = browser.new_page(viewport={"width": 1280, "height": 720})
                page.set_content(_wrap_section_page(section), wait_until="networkidle")
                pngs.append(page.screenshot(full_page=False))
            browser.close()
    except Exception:
        pngs = [img for _, img in _export_png_fallback(sections)]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for png_bytes in pngs:
        slide = prs.slides.add_slide(blank_layout)
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
            tmp.write(png_bytes)
            tmp.flush()
            slide.shapes.add_picture(tmp.name, 0, 0, width=prs.slide_width, height=prs.slide_height)
            Path(tmp.name).unlink(missing_ok=True)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _export_csv(html: str) -> bytes:
    sections = _split_sections(html)
    if not sections:
        sections = [html]

    buf = io.StringIO()
    writer = csv.writer(buf)
    writer.writerow(["page", "heading", "content"])
    for index, section in enumerate(sections, start=1):
        slide_match = re.search(r'data-slide=["\']?(\d+)["\']?', section, re.IGNORECASE)
        page = slide_match.group(1) if slide_match else str(index)
        heading_match = re.search(r"<h[1-6][^>]*>(.*?)</h[1-6]>", section, re.DOTALL | re.IGNORECASE)
        heading = _html_fragment_to_text(heading_match.group(1)) if heading_match else ""
        content = _html_fragment_to_text(section)
        writer.writerow([page, heading, content])
    return buf.getvalue().encode("utf-8")


def _html_fragment_to_text(fragment: str) -> str:
    text = re.sub(r"<style[^>]*>.*?</style>", " ", fragment, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r"<script[^>]*>.*?</script>", " ", text, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r"<[^>]+>", " ", text)
    text = html_lib.unescape(text)
    return re.sub(r"\s+", " ", text).strip()


def _split_sections(html: str) -> List[str]:
    parts = re.findall(r"<section[^>]*>.*?</section>", html, re.DOTALL | re.IGNORECASE)
    if parts:
        return parts
    articles = re.findall(r"<article[^>]*>.*?</article>", html, re.DOTALL | re.IGNORECASE)
    if articles:
        return articles
    figures = re.findall(r"<figure[^>]*>.*?</figure>", html, re.DOTALL | re.IGNORECASE)
    return figures if figures else [html]


def _wrap_section_page(section: str) -> str:
    return (
        "<!DOCTYPE html><html><head><meta charset=\"utf-8\">"
        "<style>body{margin:0;background:#fff;} section{aspect-ratio:16/9;width:1280px;height:720px;}</style>"
        f"</head><body>{section}</body></html>"
    )
