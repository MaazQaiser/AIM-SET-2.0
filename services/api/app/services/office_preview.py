from __future__ import annotations

import io
import subprocess
import tempfile
from pathlib import Path
from typing import List


OFFICE_PREVIEW_EXTENSIONS = {".ppt", ".pptx", ".docx"}
PRESENTATION_EXTENSIONS = {".ppt", ".pptx"}
DEFAULT_SLIDE_DPI = 150


def slide_storage_path(tenant_id: str, asset_id: str, slide_index: int) -> str:
    return f"{tenant_id}/{asset_id}/slides/{slide_index:03d}.png"


def convert_office_bytes_to_pdf(file_bytes: bytes, ext: str, *, allow_text_fallback: bool = False) -> bytes:
    """Convert office document bytes to PDF via LibreOffice (headless)."""
    normalized = ext.lower() if ext.startswith(".") else f".{ext.lower()}"
    if normalized not in OFFICE_PREVIEW_EXTENSIONS:
        raise ValueError(f"Unsupported office preview extension: {ext}")

    with tempfile.TemporaryDirectory() as tmp:
        work_dir = Path(tmp)
        src = work_dir / f"source{normalized}"
        src.write_bytes(file_bytes)
        pdf_path = convert_office_file_to_pdf(src, work_dir, allow_text_fallback=allow_text_fallback)
        return pdf_path.read_bytes()


def rasterize_presentation_slides(
    file_bytes: bytes,
    ext: str,
    *,
    dpi: int = DEFAULT_SLIDE_DPI,
) -> List[bytes]:
    """Render presentation slides to PNG bytes via LibreOffice -> PDF -> PNG.

    Visual slide previews require LibreOffice so uploaded decks keep their original
    layout, fonts, and branding. Text-only fallback PDFs are never used here.
    """
    normalized = ext.lower() if ext.startswith(".") else f".{ext.lower()}"
    if normalized not in PRESENTATION_EXTENSIONS:
        raise ValueError(f"Not a presentation: {ext}")

    with tempfile.TemporaryDirectory() as tmp:
        work_dir = Path(tmp)
        src = work_dir / f"source{normalized}"
        src.write_bytes(file_bytes)
        pdf_path = convert_office_file_to_pdf(src, work_dir, allow_text_fallback=False)
        return pdf_path_to_png_bytes(pdf_path, dpi=dpi)


def convert_office_file_to_pdf(
    src: Path,
    work_dir: Path,
    *,
    allow_text_fallback: bool = False,
) -> Path:
    try:
        subprocess.run(
            [
                "soffice",
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                str(work_dir),
                str(src),
            ],
            check=True,
            capture_output=True,
            timeout=180,
        )
        pdf = work_dir / f"{src.stem}.pdf"
        if pdf.is_file():
            return pdf
    except (FileNotFoundError, subprocess.CalledProcessError, subprocess.TimeoutExpired):
        pass

    if allow_text_fallback and src.suffix.lower() in PRESENTATION_EXTENSIONS:
        return _pptx_fallback_pdf(src)

    raise ValueError(
        "LibreOffice is required to render presentation previews with original design. "
        "Install LibreOffice (included in Dockerfile.api) and re-embed the asset."
    )


def pdf_path_to_png_bytes(pdf_path: Path, *, dpi: int = DEFAULT_SLIDE_DPI) -> List[bytes]:
    try:
        from pdf2image import convert_from_path  # type: ignore

        images = convert_from_path(str(pdf_path), dpi=dpi)
        out: List[bytes] = []
        for img in images:
            buf = io.BytesIO()
            img.save(buf, format="PNG")
            out.append(buf.getvalue())
        if out:
            return out
    except Exception:
        pass

    try:
        import fitz  # type: ignore

        doc = fitz.open(str(pdf_path))
        scale = dpi / 72.0
        matrix = fitz.Matrix(scale, scale)
        out = []
        for page in doc:
            pix = page.get_pixmap(matrix=matrix, alpha=False)
            out.append(pix.tobytes("png"))
        doc.close()
        if out:
            return out
    except Exception as exc:
        raise ValueError(f"PDF rasterize failed: {exc}") from exc

    raise ValueError("Could not rasterize presentation PDF to slide images")


def _pptx_fallback_pdf(src: Path) -> Path:
    """Text-only fallback for non-visual pipelines (never used for slide preview)."""
    import fitz  # type: ignore
    from pptx import Presentation  # type: ignore

    prs = Presentation(str(src))
    doc = fitz.open()
    for slide in prs.slides:
        page = doc.new_page(width=1280, height=720)
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text[:200])
        page.insert_text((40, 40), "\n".join(texts[:8]) or "Slide", fontsize=14)
    out = src.parent / f"{src.stem}_fallback.pdf"
    doc.save(str(out))
    doc.close()
    return out
