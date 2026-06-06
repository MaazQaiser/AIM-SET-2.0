from __future__ import annotations

from io import BytesIO

import pytest
from dc_core.tenancy import TenantContext

from app.config import get_settings
from app.domain.content_studio_repository import ContentStudioRepository
from app.domain.memory_store import get_memory_store


def _reset_content_template_memory() -> None:
    store = get_memory_store()
    store.content_templates.clear()
    store.content_template_html.clear()
    store.content_template_files.clear()
    store.tenant_uuid_map.clear()


def test_templates_are_public_across_tenants_in_memory_fallback(monkeypatch) -> None:
    get_settings.cache_clear()
    settings = get_settings()
    monkeypatch.setattr(settings, "supabase_url", "")
    monkeypatch.setattr(settings, "supabase_service_role_key", "")
    _reset_content_template_memory()

    repo = ContentStudioRepository()
    ctx_a = TenantContext(tenant_id="tenant-a", user_id="user-a", clerk_org_id="tenant-a")
    ctx_b = TenantContext(tenant_id="tenant-b", user_id="user-b", clerk_org_id="tenant-b")

    created = repo.create_manual_template(
        ctx_a,
        name="Shared executive deck",
        artifact_type="deck",
        html="<!doctype html><html><body><section data-slide='1'>Shared</section></body></html>",
        css_variables={"--accent": "#2563eb"},
        tags=["shared"],
        page_count=1,
    )

    shared_key = settings.kb_shared_tenant_key
    assert len(get_memory_store().content_templates.get(shared_key, [])) == 1

    listed_for_b = repo.list_templates(ctx_b, artifact_type="deck")
    assert [template["id"] for template in listed_for_b] == [created["id"]]

    loaded_for_b = repo.get_template(ctx_b, created["id"])
    assert loaded_for_b
    assert loaded_for_b["html"]

    updated = repo.update_template(ctx_b, created["id"], {"name": "Renamed public deck"})
    assert updated and updated["name"] == "Renamed public deck"
    assert repo.get_template(ctx_a, created["id"])["name"] == "Renamed public deck"

    assert repo.delete_template(ctx_b, created["id"]) is True
    assert repo.list_templates(ctx_a, artifact_type="deck") == []


def test_template_metadata_is_public_across_tenants_in_memory_fallback(monkeypatch) -> None:
    get_settings.cache_clear()
    settings = get_settings()
    monkeypatch.setattr(settings, "supabase_url", "")
    monkeypatch.setattr(settings, "supabase_service_role_key", "")
    _reset_content_template_memory()

    repo = ContentStudioRepository()
    ctx_a = TenantContext(tenant_id="tenant-a", user_id="user-a", clerk_org_id="tenant-a")
    ctx_b = TenantContext(tenant_id="tenant-b", user_id="user-b", clerk_org_id="tenant-b")

    created = repo.create_manual_template(
        ctx_a,
        name="Analyzed deck",
        artifact_type="deck",
        html="<section class='slide' data-slide='1'>Executive Snapshot</section>",
        css_variables={},
        tags=[],
        page_count=1,
    )
    repo.finalize_template(
        ctx_a,
        created["id"],
        html="<section class='slide' data-slide='1'>Executive Snapshot</section>",
        css_variables={},
        page_count=1,
        metadata={
            "slideCount": 1,
            "slides": [{"slide": 1, "title": "Executive Snapshot", "layout": "Title Slide"}],
            "design": {"colors": ["#123456"], "fonts": ["Aptos"], "layouts": ["Title Slide"]},
        },
    )

    loaded = repo.get_template(ctx_b, created["id"])
    assert loaded
    assert loaded["metadata"]["slideCount"] == 1
    assert loaded["metadata"]["slides"][0]["title"] == "Executive Snapshot"
    assert loaded["metadata"]["design"]["colors"] == ["#123456"]


def test_pptx_template_metadata_extraction_reads_title_text_and_design() -> None:
    pytest.importorskip("pptx")
    from pptx import Presentation  # type: ignore
    from pptx.dml.color import RGBColor  # type: ignore
    from pptx.enum.shapes import MSO_SHAPE  # type: ignore
    from pptx.util import Inches  # type: ignore

    from app.services.template_ingest_service import _extract_template_metadata

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "AI Sales Deck"
    box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(1))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "Pipeline proof points and customer outcomes"
    run.font.name = "Aptos"
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3), Inches(2), Inches(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(18, 52, 86)

    buf = BytesIO()
    prs.save(buf)
    metadata = _extract_template_metadata(
        file_bytes=buf.getvalue(),
        ext=".pptx",
        storage_path="tenant/template/ai-sales-deck.pptx",
        page_count=1,
        preview_image_count=1,
        has_preview_pdf=False,
        preview_pdf_bytes=None,
    )

    assert metadata["slideCount"] == 1
    assert metadata["slides"][0]["title"] == "AI Sales Deck"
    assert "Pipeline proof points" in metadata["slides"][0]["text"]
    assert "#123456" in metadata["slides"][0]["colors"]
    assert "Aptos" in metadata["design"]["fonts"]


def test_pptx_metadata_can_generate_html_fallback() -> None:
    from app.services.template_ingest_service import _metadata_css_vars, _metadata_to_template_html

    metadata = {
        "slideCount": 1,
        "slides": [
            {
                "slide": 1,
                "title": "AI Sales Deck",
                "layout": "Title Slide",
                "textBlocks": ["Pipeline proof points and customer outcomes"],
                "colors": ["#123456"],
            }
        ],
        "design": {"colors": ["#123456"], "fonts": ["Aptos"], "layouts": ["Title Slide"]},
    }

    css_vars = _metadata_css_vars(metadata)
    html = _metadata_to_template_html(metadata, css_variables=css_vars)

    assert "--template-primary" in html
    assert "AI Sales Deck" in html
    assert "Pipeline proof points and customer outcomes" in html
    assert "#123456" in html
